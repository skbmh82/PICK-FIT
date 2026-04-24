from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x0F, 0x1A)
BG_PURPLE  = RGBColor(0x2D, 0x0A, 0x5E)
BG_MID     = RGBColor(0x1A, 0x05, 0x33)
VIOLET     = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_L   = RGBColor(0xA7, 0x8B, 0xFA)
FUCHSIA    = RGBColor(0xE8, 0x79, 0xF9)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_60   = RGBColor(0x99, 0x99, 0xBB)
WHITE_40   = RGBColor(0x77, 0x77, 0x99)
CARD_BG    = RGBColor(0x1E, 0x1E, 0x30)
CARD_BG2   = RGBColor(0x28, 0x10, 0x50)
ACTIVE_CARD = RGBColor(0x3B, 0x1A, 0x7A)

blank_layout = prs.slide_layouts[6]  # completely blank

# ── 헬퍼 함수 ────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_rect(slide, l, t, w, h, fill=CARD_BG, radius=False):
    shape = slide.shapes.add_shape(
        5 if radius else 1,  # 5=rounded rect, 1=rect
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0x3A, 0x1A, 0x6A)
    shape.line.width = Pt(0.5)
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_tag(slide, text, l=0.5, t=0.4):
    add_rect(slide, l, t, 2.6, 0.38, fill=RGBColor(0x3A, 0x1A, 0x7A), radius=True)
    add_text(slide, text, l+0.1, t+0.04, 2.4, 0.35,
             size=10, bold=True, color=VIOLET_L, align=PP_ALIGN.CENTER)

def add_title(slide, line1, line2="", l=0.5, t=0.9, size1=40, size2=36):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(12), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = line1
    r.font.size = Pt(size1); r.font.bold = True; r.font.color.rgb = WHITE
    if line2:
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = line2
        r2.font.size = Pt(size2); r2.font.bold = True; r2.font.color.rgb = FUCHSIA

def add_card(slide, l, t, w, h, title, body_lines, icon="", active=False):
    fill = ACTIVE_CARD if active else CARD_BG
    add_rect(slide, l, t, w, h, fill=fill, radius=True)
    cy = t + 0.18
    if icon:
        add_text(slide, icon, l+0.15, cy, w-0.3, 0.45, size=22, align=PP_ALIGN.LEFT)
        cy += 0.45
    add_text(slide, title, l+0.15, cy, w-0.3, 0.38,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    cy += 0.40
    for line in body_lines:
        add_text(slide, f"• {line}", l+0.15, cy, w-0.3, 0.28,
                 size=9.5, color=WHITE_60, align=PP_ALIGN.LEFT)
        cy += 0.26

# ════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
add_bg(s1, BG_PURPLE)

# 배경 장식 사각형
deco = s1.shapes.add_shape(1, Inches(8), Inches(-1), Inches(6), Inches(6))
deco.fill.solid(); deco.fill.fore_color.rgb = RGBColor(0x4C, 0x1D, 0x95)
deco.line.fill.background(); deco.rotation = 30

add_tag(s1, "AI 피팅룸, 픽핏  ·  PICK FIT", l=4.8, t=1.6)

tb = s1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "픽!  한 번에 핏!"
r.font.size = Pt(60); r.font.bold = True; r.font.color.rgb = WHITE

p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "입어보고 사세요. AI로, 지금 바로."
r2.font.size = Pt(24); r2.font.bold = False; r2.font.color.rgb = FUCHSIA

add_text(s1, "Pi Network 생태계 기반 AI 가상 피팅 서비스  ·  Boutique (KANG SEONGSHIN)",
         1, 5.5, 11.3, 0.5, size=13, color=WHITE_40, align=PP_ALIGN.CENTER)

# 하단 구분선
line = s1.shapes.add_shape(1, Inches(3), Inches(5.3), Inches(7.3), Inches(0.02))
line.fill.solid(); line.fill.fore_color.rgb = VIOLET
line.line.fill.background()

# ════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_bg(s2, BG_DARK)
add_tag(s2, "PROBLEM")
add_title(s2, "한국 패션 앱에", "가상 피팅이 없다", t=0.85, size1=36, size2=32)

problems = [
    ("😤", "반품률 30% 이상", ["사이즈·핏 불일치로 의류 반품 1위", "소비자·판매자 모두 손실 발생"]),
    ("📱", "주요 플랫폼 미지원", ["무신사·29CM·지그재그·W컨셉", "AI 피팅 기능 전무 — Blue Ocean"]),
    ("💸", "충동구매 & 후회 반복", ["입어보지 못하고 구매 → 반품", "소비자 신뢰도 하락 & 환경 낭비"]),
    ("🌏", "글로벌 선점 기회", ["글로벌 경쟁사 한국 진출 전", "지금이 선점의 최적 타이밍"]),
]
xs = [0.4, 3.6, 6.8, 10.0]
for (icon, title, lines), x in zip(problems, xs):
    add_card(s2, x, 2.55, 2.9, 3.9, title, lines, icon=icon)

# ════════════════════════════════════════════════════
# SLIDE 3 — Solution
# ════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_bg(s3, BG_MID)
add_tag(s3, "SOLUTION")
add_title(s3, "PICK FIT", "AI 가상 피팅룸", t=0.85, size1=38, size2=32)

sol_box = add_rect(s3, 1.2, 2.7, 10.9, 3.8, fill=CARD_BG2, radius=True)
sol_box.line.color.rgb = VIOLET
sol_box.line.width = Pt(1)

sol_lines = [
    ("쇼핑몰 URL 하나만 붙여넣으면", WHITE_60, 14, False),
    ("내 사진 위에 옷을 AI가 직접 입혀드립니다.", WHITE, 20, True),
    ("", WHITE, 8, False),
    ("무신사 · 29CM · 지그재그 · 쿠팡 등  모든 쇼핑몰  지원", WHITE_60, 13, False),
    ("상의 · 하의 · 원피스 · 아우터  카테고리별 정밀 피팅", WHITE_60, 13, False),
    ("화면 캡쳐 → 드래그 영역 선택  · 이미지 직접 업로드", WHITE_60, 13, False),
    ("결과는  저장 · 다른 옷 입어보기  까지 한 번에", WHITE_60, 13, False),
]
y = 2.85
for (txt, col, sz, bd) in sol_lines:
    add_text(s3, txt, 1.5, y, 10.3, 0.42, size=sz, bold=bd, color=col, align=PP_ALIGN.CENTER)
    y += 0.42 if txt else 0.2

# ════════════════════════════════════════════════════
# SLIDE 4 — Features
# ════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_bg(s4, BG_DARK)
add_tag(s4, "CORE FEATURES")
add_title(s4, "MVP", "7개 핵심 화면", t=0.85, size1=36, size2=30)

feats = [
    ("01", "🔐", "로그인 · 회원가입"),
    ("02", "🏠", "홈 · 피팅 피드"),
    ("03", "✂️", "아이템 추가\nURL·업로드·캡쳐"),
    ("04", "⏳", "AI 피팅 로딩\n팁·단계 애니메이션"),
    ("05", "✨", "피팅 결과\n저장 · 다른 옷"),
    ("06", "👗", "옷장 (워드로브)\n삭제 · 필터"),
    ("07", "👤", "프로필 · 구독"),
    ("★", "🗑️", "피팅 기록 관리\n자동 삭제 Cron"),
]
col_w = 1.52
for i, (num, icon, title) in enumerate(feats):
    x = 0.3 + i * (col_w + 0.06)
    box = add_rect(s4, x, 2.55, col_w, 3.8, fill=CARD_BG if num != "★" else CARD_BG2, radius=True)
    if num == "★":
        box.line.color.rgb = VIOLET; box.line.width = Pt(1)
    add_text(s4, num, x, 2.65, col_w, 0.3, size=9, bold=True, color=VIOLET_L, align=PP_ALIGN.CENTER)
    add_text(s4, icon, x, 3.05, col_w, 0.55, size=28, align=PP_ALIGN.CENTER)
    add_text(s4, title, x, 3.65, col_w, 0.7, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 5 — Tech Stack
# ════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_bg(s5, BG_MID)
add_tag(s5, "TECHNOLOGY")
add_title(s5, "기술 스택", "", t=0.85, size1=36)

tech = [
    ("Frontend / Backend", [
        "Next.js 14  (App Router)",
        "TypeScript  Strict Mode",
        "Tailwind CSS",
        "Server Actions",
        "Vercel 배포",
    ]),
    ("Database / Auth", [
        "Supabase PostgreSQL",
        "Supabase Auth",
        "Supabase Storage",
        "Row Level Security (RLS)",
        "Admin SDK (서버 전용)",
    ]),
    ("AI · API", [
        "IDM-VTON via Replicate",
        "→ Fashn.ai 상용 전환 예정",
        "비동기 폴링 아키텍처",
        "Provider 추상화 패턴",
        "카카오 SDK (결과 공유)",
    ]),
]
for i, (title, items) in enumerate(tech):
    x = 0.4 + i * 4.3
    add_rect(s5, x, 2.35, 4.0, 4.8, fill=CARD_BG, radius=True)
    add_text(s5, title, x+0.2, 2.5, 3.6, 0.38, size=11, bold=True, color=VIOLET_L)
    y = 3.0
    for item in items:
        add_text(s5, f"  {item}", x+0.2, y, 3.6, 0.36, size=12, color=WHITE_60)
        y += 0.42

# ════════════════════════════════════════════════════
# SLIDE 6 — Business Model
# ════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_bg(s6, BG_DARK)
add_tag(s6, "BUSINESS MODEL")
add_title(s6, "요금제 구조", "", t=0.85, size1=36)

plans = [
    ("FREE",   "₩0",       "/월", "월 3회 피팅", False),
    ("BASIC",  "₩4,900",  "/월", "월 20회 피팅", True),
    ("PRO",    "₩9,900",  "/월", "월 30~50회",   False),
    ("CREDIT", "₩4,900",  "/1회", "10회 크레딧",  False),
]
for i, (name, price, unit, count, highlight) in enumerate(plans):
    x = 0.5 + i * 3.1
    fill = ACTIVE_CARD if highlight else CARD_BG
    box = add_rect(s6, x, 2.4, 2.8, 3.6, fill=fill, radius=True)
    if highlight:
        box.line.color.rgb = VIOLET; box.line.width = Pt(1.5)
    add_text(s6, name, x, 2.6, 2.8, 0.38, size=11, bold=True, color=VIOLET_L, align=PP_ALIGN.CENTER)

    tb2 = s6.shapes.add_textbox(Inches(x), Inches(3.1), Inches(2.8), Inches(0.7))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r_price = p2.add_run(); r_price.text = price
    r_price.font.size = Pt(28); r_price.font.bold = True; r_price.font.color.rgb = WHITE
    r_unit = p2.add_run(); r_unit.text = unit
    r_unit.font.size = Pt(13); r_unit.font.color.rgb = WHITE_60

    add_text(s6, count, x, 3.85, 2.8, 0.38, size=13, color=WHITE_60, align=PP_ALIGN.CENTER)
    if highlight:
        add_text(s6, "★ 인기", x, 4.3, 2.8, 0.38, size=12, bold=True, color=FUCHSIA, align=PP_ALIGN.CENTER)

add_text(s6, "Phase 3 Pi 연동 시 · Basic = 4,000 PICK (≈₩4,000, 18% 할인)",
         0.5, 6.3, 12.3, 0.45, size=12, color=WHITE_40, align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════
# SLIDE 7 — Market
# ════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_bg(s7, BG_MID)
add_tag(s7, "MARKET")
add_title(s7, "시장 규모 & 타겟", "", t=0.85, size1=36)

markets = [
    ("20조+",   "국내 온라인 패션 시장", ["2024년 기준 시장 규모", "연 10% 이상 성장 중"]),
    ("20~30대", "핵심 타겟 여성",        ["온라인 패션 구매 1위 연령대", "AI 기술 수용도 높음"]),
    ("MAU 2,500", "손익분기점",          ["유료 전환율 5% 기준", "MAU 10,000명에서 흑자"]),
    ("$0.03",   "피팅 1회 원가",         ["Kolors API 기준", "구독 모델로 마진 확보"]),
]
for i, (big, title, lines) in enumerate(markets):
    x = 0.4 + i * 3.2
    add_rect(s7, x, 2.4, 3.0, 4.2, fill=CARD_BG, radius=True)
    add_text(s7, big,   x+0.15, 2.55, 2.7, 0.65, size=26, bold=True, color=FUCHSIA)
    add_text(s7, title, x+0.15, 3.25, 2.7, 0.42, size=13, bold=True, color=WHITE)
    y = 3.75
    for line in lines:
        add_text(s7, f"• {line}", x+0.15, y, 2.7, 0.32, size=10.5, color=WHITE_60)
        y += 0.32

# ════════════════════════════════════════════════════
# SLIDE 8 — Roadmap
# ════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
add_bg(s8, BG_DARK)
add_tag(s8, "ROADMAP")
add_title(s8, "개발 로드맵", "", t=0.85, size1=36)

phases = [
    ("PHASE 1 · MVP", "현재 진행 중 🔥",
     ["이메일 로그인 + Supabase Auth", "AI 가상 피팅 (IDM-VTON)", "화면 캡쳐 · 영역 선택 기능", "옷장 · 피팅 기록 관리", "6월 런칭 목표"], True),
    ("PHASE 2 · 3~6개월", "기능 확장",
     ["친구 초대 시스템", "쇼핑몰 제휴 확대", "AI 스타일 추천", "구독 결제 연동", "Fashn.ai 전환"], False),
    ("PHASE 3 · 6~9개월", "Pi 생태계 연동",
     ["PICK 토큰 발행", "Pi Network 완전 연동", "Pi 결제 구독", "PICK PICK 공유 지갑", "Pi 유저 마이그레이션"], False),
    ("PHASE 4 · 9~18개월", "글로벌 확장",
     ["React Native 앱", "일본 · 베트남 진출", "B2B API 제공", "쇼핑몰 임베드 SDK", "IPO 준비"], False),
]
for i, (phase, title, items, active) in enumerate(phases):
    x = 0.3 + i * 3.25
    fill = ACTIVE_CARD if active else CARD_BG
    box = add_rect(s8, x, 2.4, 3.0, 4.7, fill=fill, radius=True)
    if active:
        box.line.color.rgb = VIOLET; box.line.width = Pt(1.5)
    add_text(s8, phase, x+0.15, 2.52, 2.7, 0.32, size=9, bold=True, color=VIOLET_L)
    add_text(s8, title, x+0.15, 2.88, 2.7, 0.38, size=13, bold=True, color=WHITE)
    y = 3.35
    for item in items:
        add_text(s8, f"→ {item}", x+0.15, y, 2.7, 0.32, size=10, color=WHITE_60)
        y += 0.34

# ════════════════════════════════════════════════════
# SLIDE 9 — Pi Network
# ════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
add_bg(s9, BG_MID)
add_tag(s9, "PI NETWORK STRATEGY")
add_title(s9, "Pi Network", "연동 전략", t=0.85, size1=36, size2=30)

pi_cols = [
    ("🅰️ Plan A — Pi 활성화 시", [
        "Pi Browser 우선 배포",
        "Basic = 4,000 PICK (18% 할인)",
        "Pi 커뮤니티 MAU 즉시 확보",
        "PICK PICK 배달앱과 토큰 공유",
        "Pi → PICK 전환 지갑 제공",
    ]),
    ("🅱️ Plan B — Pi 지연 시", [
        "일반 소비자 서비스로 자동 전환",
        "카카오페이 · 토스페이 결제",
        "환경변수 하나로 플랜 분기",
        "Pi 대기 중에도 서비스 지속",
        "Mainnet 개방 시 즉시 전환",
    ]),
    ("🔗 PICK PICK 시너지", [
        "배달앱 + 피팅앱 PICK 생태계",
        "공유 지갑 · 공유 회원",
        "패션 → 배달 크로스 프로모션",
        "Pi 유저베이스 공유",
        "원 플랫폼, 멀티 서비스",
    ]),
]
for i, (title, items) in enumerate(pi_cols):
    x = 0.4 + i * 4.3
    add_rect(s9, x, 2.4, 4.0, 4.8, fill=CARD_BG, radius=True)
    add_text(s9, title, x+0.2, 2.55, 3.6, 0.42, size=12, bold=True, color=WHITE)
    y = 3.1
    for item in items:
        add_text(s9, f"• {item}", x+0.2, y, 3.6, 0.35, size=11, color=WHITE_60)
        y += 0.38

# ════════════════════════════════════════════════════
# SLIDE 10 — KPI
# ════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
add_bg(s10, BG_DARK)
add_tag(s10, "NORTH STAR METRIC")
add_title(s10, "주간 피팅 완료 수", "6개월 KPI 목표", t=0.85, size1=32, size2=28)

kpis = [
    ("10,000명", "MAU 목표\n(6개월 내)"),
    ("5,000회",  "월간 피팅\n완료 수"),
    ("5%",       "유료\n전환율"),
    ("10%",      "D30\n리텐션"),
]
for i, (val, label) in enumerate(kpis):
    x = 0.5 + i * 3.1
    add_rect(s10, x, 2.5, 2.9, 3.5, fill=CARD_BG2, radius=True)
    box = s10.shapes[-1]; box.line.color.rgb = VIOLET; box.line.width = Pt(0.8)
    add_text(s10, val,   x, 3.0, 2.9, 0.8, size=32, bold=True, color=FUCHSIA, align=PP_ALIGN.CENTER)
    add_text(s10, label, x, 3.85, 2.9, 0.8, size=14, color=WHITE_60, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 11 — Close
# ════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
add_bg(s11, BG_PURPLE)

deco2 = s11.shapes.add_shape(1, Inches(9), Inches(4), Inches(5), Inches(5))
deco2.fill.solid(); deco2.fill.fore_color.rgb = RGBColor(0x5B, 0x21, 0xB6)
deco2.line.fill.background(); deco2.rotation = 25

add_tag(s11, "THANK YOU", l=5.1, t=1.3)

tb = s11.shapes.add_textbox(Inches(1), Inches(1.9), Inches(11.3), Inches(3))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "픽!  한 번에 핏!"
r.font.size = Pt(56); r.font.bold = True; r.font.color.rgb = WHITE

p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "AI 피팅룸, PICK FIT"
r2.font.size = Pt(22); r2.font.color.rgb = FUCHSIA

line2 = s11.shapes.add_shape(1, Inches(3.5), Inches(5.1), Inches(6.3), Inches(0.02))
line2.fill.solid(); line2.fill.fore_color.rgb = VIOLET
line2.line.fill.background()

add_text(s11, "Boutique (KANG SEONGSHIN)", 1, 5.3, 11.3, 0.4, size=14, color=WHITE_60, align=PP_ALIGN.CENTER)
add_text(s11, "Pi Network 생태계 기반  ·  PICK FIT + PICK PICK — 원 생태계, 멀티 서비스",
         1, 5.75, 11.3, 0.4, size=12, color=WHITE_40, align=PP_ALIGN.CENTER)

# ── 저장 ─────────────────────────────────────────────
out = r"c:\PICK FIT project\pick-fit\PICKFIT_Presentation.pptx"
prs.save(out)
print(f"저장 완료: {out}")
